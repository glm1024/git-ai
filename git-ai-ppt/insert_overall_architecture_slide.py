#!/usr/bin/env python3
"""Insert the approved dual-lane architecture slide after section 04.

The existing 39 slide parts are restored byte-for-byte after python-pptx saves
the package. The new slide remains fully editable native PowerPoint content.
"""

from __future__ import annotations

import shutil
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
SOURCE = ROOT / "git-ai-research-record.pptx"
CANDIDATE = ROOT / "git-ai-research-record.architecture-candidate.pptx"

FONT = "Microsoft YaHei"
PRIMARY_BLUE = "0062AC"
DEEP_BLUE = "00518E"
DARK_NAVY = "213261"
BODY = "111111"
MUTED = "5B6775"
LINE = "C8D7E3"
MID_BLUE = "8DC5E8"
PALE_BLUE = "EAF4FB"
PALE_NAVY = "F3F5FA"
SOFT_WHITE = "FAFAFA"
WHITE = "FFFFFF"
GREEN = "2E8B57"
PALE_GREEN = "EDF8F1"
AMBER = "D97706"
PALE_AMBER = "FFF7E8"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_run_font(run, size: float, color: str = BODY, bold: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    rpr = run._r.get_or_add_rPr()
    rpr.set(qn("a:ea"), FONT)
    rpr.set(qn("a:latin"), FONT)


def add_text(
    slide,
    content: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 13,
    color: str = BODY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing: float = 1.25,
    margin: float = 0.0,
    name: str | None = None,
    fit: bool = False,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    if fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for index, line in enumerate(content.split("\n")):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = align
        paragraph.line_spacing = line_spacing
        paragraph.space_after = Pt(0)
        run = paragraph.add_run()
        run.text = line
        set_run_font(run, size, color, bold)
    return shape


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str = LINE,
    rounded: bool = True,
    width: float = 1.0,
    name: str | None = None,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    if name:
        shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(width)
    return shape


def add_title(slide, title: str) -> None:
    # Match the established title signature in this deck.
    add_text(
        slide,
        title,
        0.36,
        0.17,
        12.20,
        0.55,
        size=28,
        color=DARK_NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
        name="页面标题",
    )


def add_section_label(slide, text_value: str, x: float, y: float, w: float, color: str) -> None:
    add_rect(slide, x, y, w, 0.27, fill=color, line=color, rounded=True, width=0)
    add_text(
        slide,
        text_value,
        x + 0.04,
        y + 0.015,
        w - 0.08,
        0.22,
        size=10.2,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )


def add_lane_label(slide, text_value: str, x: float, y: float, w: float, h: float, color: str) -> None:
    add_rect(slide, x, y, w, h, fill=color, line=color, rounded=True, width=0)
    add_text(
        slide,
        text_value,
        x + 0.08,
        y + 0.06,
        w - 0.16,
        h - 0.12,
        size=15.0,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.05,
    )


def add_stage(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str,
    *,
    fill: str,
    line: str,
    title_color: str = DARK_NAVY,
    title_size: float = 13.2,
    body_size: float = 10.7,
) -> None:
    add_rect(slide, x, y, w, h, fill=fill, line=line, rounded=True, width=1.1)
    add_text(
        slide,
        title,
        x + 0.13,
        y + 0.13,
        w - 0.26,
        0.42,
        size=title_size,
        color=title_color,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )
    add_text(
        slide,
        body,
        x + 0.13,
        y + 0.59,
        w - 0.26,
        h - 0.72,
        size=body_size,
        color=BODY,
        valign=MSO_ANCHOR.TOP,
        line_spacing=1.18,
        fit=True,
    )


def add_chevron(slide, x: float, y: float, *, color: str = MID_BLUE, w: float = 0.20, h: float = 0.30) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()


def add_merge_arrow(slide, x: float, y: float, rotation: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(x),
        Inches(y),
        Inches(0.66),
        Inches(0.17),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(PRIMARY_BLUE)
    shape.line.fill.background()
    shape.rotation = rotation


def build_architecture_slide(slide) -> None:
    add_title(slide, "双链路采集、统一入库确认：从编辑事实到正式统计")
    add_text(
        slide,
        "Kilo v5 与 Git AI 分别构造来源证据，后端保留各自归因方式，并在真实仓库确认后统一形成正式统计。",
        0.42,
        0.92,
        12.10,
        0.38,
        size=14.0,
        color=DARK_NAVY,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
        name="页面导语",
    )

    add_section_label(slide, "客户端采集与本地证据", 1.62, 1.38, 4.22, PRIMARY_BLUE)
    add_section_label(slide, "后端四重门禁：从候选事实到正式入库", 5.99, 1.38, 6.74, DEEP_BLUE)

    top_y = 1.72
    bottom_y = 3.38
    lane_h = 1.35
    box_w = 1.23
    xs = [1.62, 3.04, 4.46, 5.88, 7.30]

    add_lane_label(slide, "Kilo v5", 0.45, top_y, 0.98, lane_h, PRIMARY_BLUE)
    add_lane_label(slide, "Git AI", 0.45, bottom_y, 0.98, lane_h, DARK_NAVY)

    kilo = [
        ("编辑事实采集", "生成事实\n接受事实"),
        ("本地证据留存", "来源行账本\n失败保留"),
        ("提交证据包", "最终 Commit 差异\n待确认的 AI 行"),
        ("①可靠接收与归属", "明确回执\n身份 / 项目绑定"),
        ("②提交行归因", "完全 / 部分匹配\n处理候选变更"),
    ]
    git_ai = [
        ("多工具来源捕获", "Hook / 会话流\n工具调用"),
        ("编辑检查点", "checkpoint\n行级来源归属"),
        ("Git 证据固化", "Notes + Commit Diff\nRewrite"),
        ("①可靠接收与归属", "持久队列 / 强回执\n身份 / 项目绑定"),
        ("②来源链关联", "checkpoint + hunk\nCommit / Rewrite 关联"),
    ]

    for index, (title, body) in enumerate(kilo):
        add_stage(
            slide,
            xs[index],
            top_y,
            box_w,
            lane_h,
            title,
            body,
            fill=PALE_BLUE,
            line=MID_BLUE,
        )
        if index < len(kilo) - 1:
            add_chevron(slide, xs[index] + box_w + 0.03, top_y + 0.53)

    for index, (title, body) in enumerate(git_ai):
        add_stage(
            slide,
            xs[index],
            bottom_y,
            box_w,
            lane_h,
            title,
            body,
            fill=PALE_NAVY,
            line="9AA7CE",
        )
        if index < len(git_ai) - 1:
            add_chevron(slide, xs[index] + box_w + 0.03, bottom_y + 0.53)

    common_xs = [9.17, 10.41, 11.65]
    common = [
        ("③真实入库确认", "扫描实际仓库\n确认进入目标分支"),
        ("④逐行证据准入", "核验目标 Diff\n同一来源内唯一计数"),
        ("正式统计与展示", "committed 日事实\n趋势 / 明细 / 看板"),
    ]
    add_rect(slide, 8.62, 1.79, 3.98, 0.42, fill=PALE_AMBER, line="E7B55E", rounded=True, width=0.8)
    add_text(
        slide,
        "客户端 Commit / committed = 本地提交证据，不等于正式入库",
        8.78,
        1.86,
        3.66,
        0.26,
        size=10.2,
        color=AMBER,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
        fit=True,
    )
    add_merge_arrow(slide, 8.54, 2.64, 19)
    add_merge_arrow(slide, 8.54, 3.75, -19)
    for index, (title, body) in enumerate(common):
        add_stage(
            slide,
            common_xs[index],
            2.50,
            1.08,
            1.45,
            title,
            body,
            fill=PALE_GREEN if index == 2 else SOFT_WHITE,
            line="76B68F" if index == 2 else PRIMARY_BLUE,
            title_color=GREEN if index == 2 else DARK_NAVY,
            title_size=12.8,
            body_size=10.3,
        )
        if index < len(common) - 1:
            add_chevron(slide, common_xs[index] + 1.10, 3.07, color=PRIMARY_BLUE, w=0.17, h=0.30)

    # Token is deliberately drawn as a bypass: it does not enter code attribution.
    add_rect(slide, 0.55, 5.08, 12.00, 0.82, fill=PALE_AMBER, line="E7B55E", rounded=True, width=0.9)
    add_text(
        slide,
        "Token 旁路",
        0.78,
        5.22,
        1.35,
        0.28,
        size=12.2,
        color=AMBER,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )
    add_text(
        slide,
        "模型 / 会话用量",
        2.25,
        5.22,
        1.45,
        0.28,
        size=11.7,
        color=BODY,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )
    add_chevron(slide, 3.78, 5.22, color=AMBER, w=0.18, h=0.26)
    add_text(
        slide,
        "按日累计用量",
        4.05,
        5.22,
        1.45,
        0.28,
        size=11.7,
        color=BODY,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )
    add_chevron(slide, 5.58, 5.22, color=AMBER, w=0.18, h=0.26)
    add_text(
        slide,
        "独立接收与幂等覆盖",
        5.85,
        5.22,
        1.78,
        0.28,
        size=11.7,
        color=BODY,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )
    add_chevron(slide, 7.72, 5.22, color=AMBER, w=0.18, h=0.26)
    add_text(
        slide,
        "Token 日统计与看板",
        7.99,
        5.22,
        1.72,
        0.28,
        size=11.7,
        color=BODY,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )
    add_text(
        slide,
        "Kilo v5 与 Git AI 均按日统计；不经过 Commit 归因、目标分支确认和最终行去重。",
        9.92,
        5.16,
        2.35,
        0.46,
        size=9.7,
        color=BODY,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.15,
        fit=True,
    )

    add_rect(slide, 0.55, 6.08, 10.82, 0.74, fill=PALE_BLUE, line=MID_BLUE, rounded=True, width=0.9)
    add_text(
        slide,
        "客户端 Commit / committed 仅表示本地提交证据；后端确认代码真实进入目标分支、核验目标 Diff 并完成同一来源内唯一计数后，才进入正式统计。Token 按日独立统计，不参与代码入库归因。",
        0.78,
        6.14,
        10.36,
        0.39,
        size=11.2,
        color=DARK_NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.1,
        fit=True,
    )
    add_text(
        slide,
        "边界：同一来源内完成幂等与最终行去重；Kilo v5 与 Git AI 若同时认领同一最终行，当前不承诺跨来源逐行去重。",
        0.78,
        6.52,
        10.36,
        0.20,
        size=8.7,
        color=DARK_NAVY,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )


def restore_protected_slides(source: Path, output: Path, count: int) -> None:
    with tempfile.TemporaryDirectory(prefix="git-ai-architecture-slide-") as temp_name:
        temp = Path(temp_name)
        source_dir = temp / "source"
        output_dir = temp / "output"
        with zipfile.ZipFile(source) as archive:
            archive.extractall(source_dir)
        with zipfile.ZipFile(output) as archive:
            archive.extractall(output_dir)

        for number in range(1, count + 1):
            for relative_path in (
                Path(f"ppt/slides/slide{number}.xml"),
                Path(f"ppt/slides/_rels/slide{number}.xml.rels"),
            ):
                original = source_dir / relative_path
                generated = output_dir / relative_path
                if original.exists():
                    generated.parent.mkdir(parents=True, exist_ok=True)
                    shutil.copy2(original, generated)

        rebuilt = temp / "rebuilt.pptx"
        with zipfile.ZipFile(rebuilt, "w", zipfile.ZIP_DEFLATED) as archive:
            for file_path in sorted(output_dir.rglob("*")):
                if file_path.is_file():
                    archive.write(file_path, file_path.relative_to(output_dir).as_posix())
        shutil.copy2(rebuilt, output)


def main() -> None:
    presentation = Presentation(SOURCE)
    original_slides = list(presentation.slides)
    protected_count = len(original_slides)
    if protected_count != 39:
        raise RuntimeError(f"expected 39 source slides, got {protected_count}")

    new_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build_architecture_slide(new_slide)

    # Insert after section 04 (slide 18) and before the existing backend detail.
    desired = original_slides[:18] + [new_slide] + original_slides[18:]
    current = list(presentation.slides)
    id_map = {
        id(slide): slide_id
        for slide, slide_id in zip(current, list(presentation.slides._sldIdLst))
    }
    for slide_id in list(presentation.slides._sldIdLst):
        presentation.slides._sldIdLst.remove(slide_id)
    for slide in desired:
        presentation.slides._sldIdLst.append(id_map[id(slide)])

    presentation.save(CANDIDATE)
    restore_protected_slides(SOURCE, CANDIDATE, protected_count)
    print(f"saved: {CANDIDATE}")
    print(f"slides: {len(desired)}")
    print("inserted_at: 19")


if __name__ == "__main__":
    main()
