#!/usr/bin/env python3
"""Add one evidence-led problem slide to the current Git AI research deck.

The existing 24 slides are treated as protected content.  The script adds one
new slide, reorders only the presentation slide list, then restores every
protected slide XML part byte-for-byte from the source deck.
"""

from __future__ import annotations

import shutil
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
SOURCE = ROOT / "git-ai-research-record-updated.pptx"
OUTPUT = ROOT / "git-ai-research-record-problem-1.pptx"

FONT = "Microsoft YaHei"
PRIMARY_BLUE = "0062AC"
DEEP_BLUE = "00518E"
DARK_NAVY = "213261"
BODY = "111111"
RULE = "A4A3A4"
LIGHT_FILL = "F2F4F7"
PALE_BLUE = "EAF4FB"
MID_BLUE = "8DC5E8"
RISK_RED = "C00000"
WHITE = "FFFFFF"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_run_font(run, size: float, color: str = BODY, bold: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    rpr = run._r.get_or_add_rPr()
    rpr.set(qn("a:ea"), FONT)
    rpr.set(qn("a:latin"), FONT)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 14,
    color: str = BODY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing: float = 1.5,
    margin: float = 0.0,
    name: str | None = None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        set_run_font(run, size, color, bold)
    return shape


def add_title(slide, title: str) -> None:
    # Canonical V1 normal-content title signature from idtpptx.
    add_text(
        slide,
        title,
        0.37,
        0.00,
        12.60,
        0.96,
        size=28,
        color=DARK_NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
        name="页面标题",
    )


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str = RULE,
    rounded: bool = True,
    width: float = 1.0,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(width)
    return shape


def add_rail_block(slide, y: float, label: str, body: str, *, height: float, accent: str = PRIMARY_BLUE) -> None:
    x, w = 8.02, 4.30
    add_rect(slide, x, y, w, height, fill=WHITE, line=MID_BLUE, rounded=True, width=0.9)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(height))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = rgb(accent)
    stripe.line.fill.background()
    add_text(slide, label, x + 0.22, y + 0.13, 0.92, 0.30, size=13.5, color=accent, bold=True,
             valign=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    add_text(slide, body, x + 1.12, y + 0.10, w - 1.34, height - 0.20, size=11.6, color=BODY,
             valign=MSO_ANCHOR.MIDDLE, line_spacing=1.28)


def add_rank_bar(slide, y: float, label: str, value_text: str, ratio: float, *, accent: str) -> None:
    add_text(slide, label, 0.82, y, 1.85, 0.30, size=12.2, color=BODY, bold=True,
             valign=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    add_rect(slide, 2.72, y + 0.05, 3.20, 0.20, fill=LIGHT_FILL, line=LIGHT_FILL, rounded=False, width=0)
    width = max(0.16, 3.20 * ratio)
    add_rect(slide, 2.72, y + 0.05, width, 0.20, fill=accent, line=accent, rounded=False, width=0)
    add_text(slide, value_text, 6.02, y, 1.12, 0.30, size=12.2, color=accent, bold=True,
             align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.0)


def build_problem_slide(slide) -> None:
    add_title(slide, "问题一：巨量导入抬高分母，AI 贡献率被稀释")
    add_text(
        slide,
        "当前分母混入日常开发、整体迁移、模板基线和开源导入，项目间 AI 贡献率暂不宜直接横向比较。",
        0.53,
        1.02,
        11.85,
        0.38,
        size=14.2,
        color=DARK_NAVY,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
        name="页面导语",
    )

    # Evidence field (EVD-01 dominant proof area).
    add_rect(slide, 0.55, 1.50, 7.18, 4.74, fill=WHITE, line=MID_BLUE, rounded=True, width=1.1)
    add_text(slide, "本地 6—7 月去重样本", 0.82, 1.73, 3.10, 0.30, size=15.2, color=DARK_NAVY,
             bold=True, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    add_text(slide, "442.9 万行", 0.82, 2.10, 2.60, 0.58, size=31, color=RISK_RED, bold=True,
             valign=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    add_text(slide, "单次最大提交（新增 + 删除）", 3.18, 2.22, 2.24, 0.34, size=12.4, color=BODY,
             valign=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    add_text(slide, "icfs · 2026-07-08", 5.55, 2.22, 1.63, 0.34, size=11.0, color=DEEP_BLUE,
             bold=True, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    add_text(slide, "15 万行以上的异常候选", 0.82, 2.88, 2.38, 0.28, size=12.4, color=BODY,
             bold=True, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    add_text(slide, "8 次提交", 3.30, 2.88, 0.90, 0.28, size=12.4, color=RISK_RED, bold=True,
             valign=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    add_text(slide, "占总提交行数 79.4%", 4.30, 2.88, 2.86, 0.28, size=13.0, color=RISK_RED,
             bold=True, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    # 79.4 / 20.6 composition bar.
    add_rect(slide, 0.82, 3.24, 6.36, 0.35, fill=LIGHT_FILL, line=LIGHT_FILL, rounded=False, width=0)
    add_rect(slide, 0.82, 3.24, 5.05, 0.35, fill=RISK_RED, line=RISK_RED, rounded=False, width=0)
    add_text(slide, "异常候选 1,278.9 万行", 0.98, 3.25, 2.42, 0.30, size=10.9, color=WHITE, bold=True,
             valign=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

    add_text(slide, "最大四次提交", 0.82, 3.85, 1.60, 0.28, size=13.0, color=DARK_NAVY, bold=True,
             valign=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    add_rank_bar(slide, 4.20, "icfs · 07-08", "442.9 万", 1.0, accent=RISK_RED)
    add_rank_bar(slide, 4.61, "icfs · 06-04", "442.7 万", 442.7164 / 442.8660, accent=RISK_RED)
    add_rank_bar(slide, 5.02, "iagent · 07-20", "96.3 万", 96.3145 / 442.8660, accent=PRIMARY_BLUE)
    add_rank_bar(slide, 5.43, "iagent · 06-10", "90.1 万", 90.0906 / 442.8660, accent=PRIMARY_BLUE)
    add_text(
        slide,
        "口径：本地库 2026-06-01—07-21；按规范仓库 + commit 去重；提交行数 = 新增 + 删除。",
        0.82,
        5.84,
        6.35,
        0.28,
        size=8.7,
        color=BODY,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )

    # Interpretation rail.
    add_rail_block(
        slide,
        1.50,
        "问题",
        "少数超大提交决定了绝大部分分母，掩盖日常开发中的 AI 使用。",
        height=0.88,
        accent=RISK_RED,
    )
    add_rail_block(
        slide,
        2.52,
        "已知来源",
        "iagent：外网开发后整体导入内网\nopenspec-ai：模板 / 基线代码复制\n开源项目整体引入：项目待确认",
        height=1.34,
        accent=PRIMARY_BLUE,
    )
    add_rail_block(
        slide,
        4.00,
        "现状",
        "“扣减标注”已可临时修正首页，但仍依赖管理员填写固定行数，尚未形成来源识别与统一调整口径。",
        height=1.02,
        accent=DEEP_BLUE,
    )
    add_rail_block(
        slide,
        5.16,
        "建议",
        "系统识别异常候选 → 仓库维护人确认来源 / 排除范围 → 原始口径与调整口径并列展示。",
        height=1.08,
        accent=PRIMARY_BLUE,
    )

    # Keep the other major factor visible without diluting the slide's focus.
    add_rect(slide, 0.55, 6.38, 10.48, 0.56, fill=PALE_BLUE, line=MID_BLUE, rounded=True, width=0.9)
    add_text(
        slide,
        "另一因素是分子漏采：Codex、Claude Code、Kilo v7、OpenCode 等历史生成事实尚未完整进入统计；Git AI 上线后需单独观测。本页只聚焦分母问题。",
        0.78,
        6.45,
        10.02,
        0.40,
        size=11.2,
        color=DARK_NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.05,
    )


def restore_protected_slides(source: Path, output: Path, count: int) -> None:
    with tempfile.TemporaryDirectory(prefix="git-ai-problem-slide-") as tmpdir:
        temp = Path(tmpdir)
        src_dir = temp / "source"
        out_dir = temp / "output"
        with zipfile.ZipFile(source) as zf:
            zf.extractall(src_dir)
        with zipfile.ZipFile(output) as zf:
            zf.extractall(out_dir)

        for number in range(1, count + 1):
            for rel_path in (
                Path(f"ppt/slides/slide{number}.xml"),
                Path(f"ppt/slides/_rels/slide{number}.xml.rels"),
            ):
                src = src_dir / rel_path
                dst = out_dir / rel_path
                if src.exists():
                    dst.parent.mkdir(parents=True, exist_ok=True)
                    shutil.copy2(src, dst)

        rebuilt = temp / "rebuilt.pptx"
        with zipfile.ZipFile(rebuilt, "w", zipfile.ZIP_DEFLATED) as zf:
            for path in sorted(out_dir.rglob("*")):
                if path.is_file():
                    zf.write(path, path.relative_to(out_dir).as_posix())
        shutil.copy2(rebuilt, output)


def main() -> None:
    prs = Presentation(SOURCE)
    original_slides = list(prs.slides)
    if len(original_slides) != 24:
        raise RuntimeError(f"expected 24 source slides, got {len(original_slides)}")

    new_slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_problem_slide(new_slide)

    # Place the problem after the target-commit solution and before section 05.
    desired = original_slides[:22] + [new_slide] + original_slides[22:]
    current = list(prs.slides)
    id_map = {id(slide): sld_id for slide, sld_id in zip(current, list(prs.slides._sldIdLst))}
    for sld_id in list(prs.slides._sldIdLst):
        prs.slides._sldIdLst.remove(sld_id)
    for slide in desired:
        prs.slides._sldIdLst.append(id_map[id(slide)])

    prs.save(OUTPUT)
    restore_protected_slides(SOURCE, OUTPUT, count=24)
    print(f"saved: {OUTPUT}")
    print(f"slides: {len(desired)}")


if __name__ == "__main__":
    main()
