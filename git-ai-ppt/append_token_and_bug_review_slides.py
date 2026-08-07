#!/usr/bin/env python3
"""Append two review slides while preserving the existing 37 slides byte-for-byte."""

from __future__ import annotations

import shutil
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
SOURCE = ROOT / "git-ai-research-record-problem-1.pptx"
OUTPUT = ROOT / "git-ai-research-record-problem-1.candidate.pptx"

FONT = "Microsoft YaHei"
PRIMARY_BLUE = "0062AC"
DEEP_BLUE = "00518E"
DARK_NAVY = "213261"
BODY = "111111"
RULE = "A4A3A4"
LIGHT_RULE = "C8D7E3"
LIGHT_FILL = "F2F4F7"
PALE_BLUE = "EAF4FB"
SOFT_WHITE = "FAFAFA"
WHITE = "FFFFFF"
RISK_RED = "C00000"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_run_font(run, size: float, color: str = BODY, bold: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    rpr = run._r.get_or_add_rPr()
    rpr.set(qn("a:ea"), FONT)
    rpr.set(qn("a:latin"), FONT)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 14,
    color: str = BODY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing: float = 1.5,
    margin: float = 0.0,
    name: str | None = None,
    fit: bool = False,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    if fit:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for index, line in enumerate(text.split("\n")):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = align
        paragraph.line_spacing = line_spacing
        paragraph.space_after = Pt(0)
        run = paragraph.add_run()
        run.text = line
        set_run_font(run, size, color, bold)
    return shape


def add_title(slide, title: str) -> None:
    add_text(
        slide,
        title,
        0.37,
        0.00,
        12.60,
        0.96,
        size=28,
        color=DARK_NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
        name="页面标题",
    )


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str = LIGHT_RULE,
    rounded: bool = True,
    width: float = 1.0,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(width)
    return shape


def add_tag(slide, text: str, x: float, y: float, w: float, *, custom: bool) -> None:
    fill = PRIMARY_BLUE if custom else DARK_NAVY
    add_rect(slide, x, y, w, 0.34, fill=fill, line=fill, rounded=True, width=0)
    add_text(
        slide,
        text,
        x + 0.04,
        y + 0.02,
        w - 0.08,
        0.28,
        size=11.3,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )


def add_pipeline_stage(
    slide,
    x: float,
    title: str,
    body: str,
    *,
    custom: bool,
) -> None:
    y, w, h = 1.62, 2.18, 3.48
    fill = PALE_BLUE if custom else WHITE
    accent = PRIMARY_BLUE if custom else DARK_NAVY
    add_rect(slide, x, y, w, h, fill=fill, line=accent, rounded=True, width=1.2)
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(0.08),
        Inches(h),
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = rgb(accent)
    stripe.line.fill.background()
    add_tag(slide, "企业二开" if custom else "Git AI 原生", x + 0.20, y + 0.20, 1.10, custom=custom)
    add_text(
        slide,
        title,
        x + 0.20,
        y + 0.72,
        w - 0.40,
        0.50,
        size=17.0,
        color=DARK_NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )
    add_text(
        slide,
        body,
        x + 0.20,
        y + 1.32,
        w - 0.40,
        h - 1.56,
        size=12.4,
        color=BODY,
        valign=MSO_ANCHOR.TOP,
        line_spacing=1.35,
        fit=True,
    )


def add_chevron(slide, x: float, y: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        Inches(x),
        Inches(y),
        Inches(0.27),
        Inches(0.42),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("8DC5E8")
    shape.line.fill.background()


def build_token_slide(slide) -> None:
    add_title(slide, "Token 统计链路：复用 Git AI 通用框架，补齐企业计量闭环")
    add_text(
        slide,
        "Git AI 原有 daemon 负责触发、增量扫描、可靠队列和批量上传；企业二开只补充 Token 事实与后台统计。",
        0.53,
        1.02,
        12.00,
        0.36,
        size=14.2,
        color=DARK_NAVY,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
        name="页面导语",
    )

    stages = [
        (
            "触发扫描",
            "Codex 工具 Hook\n\n守护进程启动\n\ncommit / push\n\n每 30 分钟兜底扫描",
            False,
        ),
        (
            "增量读取",
            "会话发现\n\n按水位读取新增记录\n\n去重与断点续扫\n\n不直接调用 AI 接口",
            False,
        ),
        (
            "Token 计量",
            "解析累计 Token\n\n计算相对上次的增量\n\n按用户 / 日期聚合\n\n生成当天累计快照",
            True,
        ),
        (
            "可靠传输",
            "先写本地指标队列\n\n约每 3 秒检查一次\n\n有数据才批量上传\n\n失败保留并退避重试",
            False,
        ),
        (
            "后台落地",
            "幂等接收与版本覆盖\n\n用户 / 组织维度投影\n\n日汇总与请求数\n\n热力图、趋势和明细",
            True,
        ),
    ]
    xs = [0.48, 2.99, 5.50, 8.01, 10.52]
    for index, (title, body, custom) in enumerate(stages):
        add_pipeline_stage(slide, xs[index], title, body, custom=custom)
        if index < len(stages) - 1:
            add_chevron(slide, xs[index] + 2.25, 3.13)

    add_rect(slide, 0.55, 5.48, 11.78, 1.12, fill=SOFT_WHITE, line=LIGHT_RULE, rounded=True, width=0.9)
    add_text(
        slide,
        "产生待上传数据",
        0.82,
        5.70,
        1.58,
        0.30,
        size=14.0,
        color=DEEP_BLUE,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )
    add_text(
        slide,
        "Hook 或扫描发现新的 Token 累计值，完成增量计算并写入本地队列。",
        2.42,
        5.60,
        3.34,
        0.52,
        size=12.6,
        color=BODY,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.25,
    )
    add_chevron(slide, 5.91, 5.82)
    add_text(
        slide,
        "发送到后端",
        6.42,
        5.70,
        1.35,
        0.30,
        size=14.0,
        color=DEEP_BLUE,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )
    add_text(
        slide,
        "3 秒任务只检查并发送队列：没有新快照就不发请求，也不会重新扫描会话。",
        7.86,
        5.60,
        4.06,
        0.52,
        size=12.6,
        color=BODY,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.25,
    )


def add_risk_strip(slide, y: float, title: str, body: str) -> None:
    x, w, h = 0.55, 3.00, 0.83
    add_rect(slide, x, y, w, h, fill=WHITE, line=LIGHT_RULE, rounded=True, width=0.8)
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(0.07),
        Inches(h),
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = rgb(PRIMARY_BLUE)
    stripe.line.fill.background()
    add_text(
        slide,
        title,
        x + 0.20,
        y + 0.12,
        0.92,
        0.26,
        size=13.0,
        color=DEEP_BLUE,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )
    add_text(
        slide,
        body,
        x + 1.10,
        y + 0.08,
        1.68,
        h - 0.16,
        size=10.6,
        color=BODY,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.22,
        fit=True,
    )


def add_bug_table(slide) -> None:
    x, y, w, h = 3.80, 1.76, 8.53, 4.72
    headers = ["代表性问题", "必须出现的真实场景", "支撑"]
    rows = [
        ("同一仓库 29 行变 58 行", "Windows、macOS 路径与历史配置并存，仓库被重复识别", "强"),
        ("amend 后旧、新提交同时生效", "特定监听时序下，旧提交未被生命周期事件作废", "强"),
        ("修改 8 行却显示生成 11、入库 9", "首尾空行、删除行与上报行数组合出现", "较强"),
        ("入库率超过 100%", "重复文本、重复认领与最终内容不一致；865 行、10 个提交", "强"),
        ("组织筛选漏掉 2,226 行", "开发组织、模块归属与多种排除规则交叉重叠", "强"),
        ("“本月”筛选却显示“当天”", "用户快速切换条件，网络响应发生乱序", "中等"),
        ("事件跳过、重复或旧数据覆盖", "写入失败、部分成功、同秒碰撞、离线积压后乱序", "强（含模拟）"),
        ("约 18.7 万条后查询明显变慢", "开发数据量过小，无法体现真实查询路径瓶颈", "强"),
    ]
    table = slide.shapes.add_table(len(rows) + 1, 3, Inches(x), Inches(y), Inches(w), Inches(h)).table
    widths = [2.48, 4.95, 1.10]
    for index, width in enumerate(widths):
        table.columns[index].width = Inches(width)
    table.rows[0].height = Inches(0.44)
    for row_index in range(1, len(table.rows)):
        table.rows[row_index].height = Inches(0.535)

    for column, value in enumerate(headers):
        cell = table.cell(0, column)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(PRIMARY_BLUE)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.09)
        cell.margin_top = cell.margin_bottom = Inches(0.03)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = 1.0
        run = paragraph.add_run()
        run.text = value
        set_run_font(run, 12.2, WHITE, True)

    for row_index, row_values in enumerate(rows, start=1):
        for column, value in enumerate(row_values):
            cell = table.cell(row_index, column)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(WHITE if row_index % 2 else PALE_BLUE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.09)
            cell.margin_top = cell.margin_bottom = Inches(0.025)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.line_spacing = 1.16
            run = paragraph.add_run()
            run.text = value
            color = RISK_RED if column == 2 and value.startswith("强") else BODY
            set_run_font(run, 10.4, color, column == 2)


def build_bug_review_slide(slide) -> None:
    add_title(slide, "上线后 Bug 回顾：真实运行暴露的是组合风险")
    add_rect(slide, 0.55, 1.00, 11.78, 0.58, fill=PALE_BLUE, line="8DC5E8", rounded=True, width=0.9)
    add_text(
        slide,
        "结论：证据较强。但应表述为“真实业务数据与运行时序集中暴露了开发阶段覆盖不足”，而不是“这些 Bug 原理上无法提前发现”。",
        0.78,
        1.10,
        11.32,
        0.36,
        size=13.5,
        color=DARK_NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
        fit=True,
    )

    add_text(
        slide,
        "五类组合风险",
        0.72,
        1.78,
        2.60,
        0.32,
        size=15.0,
        color=DARK_NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )
    risks = [
        ("数据异构", "跨系统路径、历史配置、重复维度"),
        ("状态积累", "旧客户端、历史记录、新旧逻辑并存"),
        ("事件时序", "amend、重试、离线、乱序、部分成功"),
        ("组合复杂度", "单条件正常，多条件叠加才出错"),
        ("生产规模", "小样本正确，规模增长后性能失效"),
    ]
    for index, (title, body) in enumerate(risks):
        add_risk_strip(slide, 2.14 + index * 0.87, title, body)

    add_bug_table(slide)

    # Keep the conclusion strip clear of the deck's fixed bottom-right logo.
    add_rect(slide, 0.55, 6.53, 10.62, 0.47, fill=SOFT_WHITE, line=RULE, rounded=True, width=0.9)
    add_text(
        slide,
        "对外建议：系统主要风险不在单点功能，而在跨版本、跨组织、历史状态、异常时序和数据规模的组合；质量建设应升级为场景回放、不变量校验、失败注入、时间边界与容量测试。边界：当前足以定性论证，尚不能给出 Bug 占比。",
        0.78,
        6.56,
        10.14,
        0.36,
        size=11.2,
        color=DARK_NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        line_spacing=1.12,
        fit=True,
    )


def restore_protected_slides(source: Path, output: Path, count: int) -> None:
    with tempfile.TemporaryDirectory(prefix="git-ai-two-slides-") as tmpdir:
        temp = Path(tmpdir)
        source_dir = temp / "source"
        output_dir = temp / "output"
        with zipfile.ZipFile(source) as archive:
            archive.extractall(source_dir)
        with zipfile.ZipFile(output) as archive:
            archive.extractall(output_dir)

        for number in range(1, count + 1):
            for relative_path in (
                Path(f"ppt/slides/slide{number}.xml"),
                Path(f"ppt/slides/_rels/slide{number}.xml.rels"),
            ):
                original = source_dir / relative_path
                generated = output_dir / relative_path
                if original.exists():
                    generated.parent.mkdir(parents=True, exist_ok=True)
                    shutil.copy2(original, generated)

        rebuilt = temp / "rebuilt.pptx"
        with zipfile.ZipFile(rebuilt, "w", zipfile.ZIP_DEFLATED) as archive:
            for path in sorted(output_dir.rglob("*")):
                if path.is_file():
                    archive.write(path, path.relative_to(output_dir).as_posix())
        shutil.copy2(rebuilt, output)


def main() -> None:
    presentation = Presentation(SOURCE)
    protected_count = len(presentation.slides)
    if protected_count != 37:
        raise RuntimeError(f"expected 37 source slides, got {protected_count}")

    blank = presentation.slide_layouts[6]
    token_slide = presentation.slides.add_slide(blank)
    build_token_slide(token_slide)
    bug_slide = presentation.slides.add_slide(blank)
    build_bug_review_slide(bug_slide)

    presentation.save(OUTPUT)
    restore_protected_slides(SOURCE, OUTPUT, protected_count)
    print(f"saved: {OUTPUT}")
    print(f"slides: {len(presentation.slides)}")


if __name__ == "__main__":
    main()
