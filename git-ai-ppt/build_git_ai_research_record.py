#!/usr/bin/env python3
"""Update selected slides and insert new chapter slides without changing protected pages."""

from __future__ import annotations

import shutil
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
SOURCE = ROOT / "git-ai-research-record.pptx"
OUTPUT = ROOT / "git-ai-research-record-updated.pptx"

FONT = "Microsoft YaHei"
BLUE = "0068B5"
DEEP_BLUE = "004A84"
NAVY = "18345D"
LIGHT_BLUE = "EAF4FB"
PALE_BLUE = "F5FAFD"
MID_BLUE = "8DC5E8"
TEXT = "1F1F1F"
MUTED = "515151"
LINE = "C8D7E3"
RED = "C00000"
GREEN = "267A3E"
AMBER = "A85B00"
WHITE = "FFFFFF"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_run_font(run, size: float, color: str = TEXT, bold: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    rpr = run._r.get_or_add_rPr()
    rpr.set(qn("a:ea"), FONT)
    rpr.set(qn("a:latin"), FONT)


def clear_text_frame(tf) -> None:
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)


def add_text(slide, text, x, y, w, h, *, size=15, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line_spacing=1.18,
             margin=0.03, name=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = line
        set_run_font(r, size, color, bold)
    return shape


def set_existing_text(shape, text, *, size, color, bold=False, align=PP_ALIGN.LEFT,
                      valign=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    clear_text_frame(tf)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run_font(r, size, color, bold)


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.36, 0.17, 10.75, 0.55,
             size=28 if len(title) <= 25 else 26, color=NAVY, bold=True,
             valign=MSO_ANCHOR.MIDDLE, name="页面标题")
    if subtitle:
        add_text(slide, subtitle, 0.42, 0.92, 11.65, 0.42, size=15,
                 color=NAVY, valign=MSO_ANCHOR.MIDDLE, name="页面导语")


def add_rect(slide, x, y, w, h, *, fill=WHITE, line=LINE, rounded=True, width=1.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(width)
    return shape


def add_card(slide, x, y, w, h, title, body, *, header=BLUE, fill=WHITE,
             title_size=16, body_size=13.2):
    add_rect(slide, x, y, w, h, fill=fill, line=LINE, rounded=True)
    head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.48))
    head.fill.solid()
    head.fill.fore_color.rgb = rgb(header)
    head.line.fill.background()
    add_text(slide, title, x + 0.16, y + 0.06, w - 0.32, 0.34,
             size=title_size, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, body, x + 0.18, y + 0.62, w - 0.36, h - 0.76,
             size=body_size, color=TEXT, line_spacing=1.27)


def add_label(slide, text, x, y, w, *, fill=LIGHT_BLUE, color=DEEP_BLUE):
    add_rect(slide, x, y, w, 0.34, fill=fill, line=fill, rounded=True)
    add_text(slide, text, x + 0.05, y + 0.03, w - 0.10, 0.25,
             size=11.5, color=color, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)


def add_step(slide, number, title, body, x, y, w, h, *, accent=BLUE):
    add_rect(slide, x, y, w, h, fill=WHITE, line=LINE, rounded=True)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.16), Inches(y + 0.16), Inches(0.42), Inches(0.42))
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb(accent)
    circle.line.fill.background()
    add_text(slide, str(number), x + 0.16, y + 0.19, 0.42, 0.31,
             size=13.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, title, x + 0.68, y + 0.13, w - 0.82, 0.36,
             size=14.3, color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, body, x + 0.18, y + 0.68, w - 0.36, h - 0.82,
             size=12.2, color=TEXT, line_spacing=1.23)


def add_chevron(slide, x, y, w=0.27, h=0.33, color=MID_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()


def add_section(slide, chapter, title, description):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(2.05), Inches(12.05), Inches(2.20))
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(BLUE)
    band.line.fill.background()
    add_text(slide, chapter, 0.95, 2.34, 1.15, 0.50, size=28, color=WHITE,
             bold=True, valign=MSO_ANCHOR.MIDDLE)
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.05), Inches(2.33), Inches(0.03), Inches(1.45))
    divider.fill.solid()
    divider.fill.fore_color.rgb = rgb(WHITE)
    divider.line.fill.background()
    add_text(slide, title, 2.35, 2.30, 8.90, 0.72, size=32, color=WHITE,
             bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, description, 2.37, 3.16, 9.30, 0.55, size=16,
             color=WHITE, valign=MSO_ANCHOR.MIDDLE)


def add_table(slide, x, y, w, h, headers, rows, col_widths, font_size=12.2):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for i, width in enumerate(col_widths):
        table.columns[i].width = Inches(width)
    for c, value in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(BLUE)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.09)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = value
        set_run_font(r, 13, WHITE, True)
    for r_idx, row in enumerate(rows, 1):
        for c, value in enumerate(row):
            cell = table.cell(r_idx, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(WHITE if r_idx % 2 else PALE_BLUE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.09)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.18
            r = p.add_run()
            r.text = value
            set_run_font(r, font_size, TEXT, False)


def redesign_problem_slide(slide):
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)
    add_title(
        slide,
        "技术问题：本地 commit 不等于最终入库事实",
        "旧口径从本地动作或本地 commit 出发，无法确认代码是否、何时、以哪个最终 commit 进入目标仓库。",
    )

    # 场景一
    add_rect(slide, 0.55, 1.48, 5.85, 4.28, fill=WHITE, line=LINE, rounded=True)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.48), Inches(5.85), Inches(0.50))
    header.fill.solid(); header.fill.fore_color.rgb = rgb(BLUE); header.line.fill.background()
    add_text(slide, "场景一｜本地分支最终未入库", 0.76, 1.56, 5.40, 0.32,
             size=17, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    items = [
        ("7 月 1 日", "AI 生成 100 行，开发者接受 60 行"),
        ("本地 commit A", "只保留在个人分支，随后被放弃"),
        ("目标仓库", "最终入库 0 行"),
    ]
    for i, (left, right) in enumerate(items):
        y = 2.25 + i * 0.93
        add_label(slide, left, 0.82, y, 1.36)
        add_text(slide, right, 2.40, y - 0.02, 3.55, 0.48,
                 size=13.2, valign=MSO_ANCHOR.MIDDLE)
        if i < 2:
            add_chevron(slide, 3.35, y + 0.57, w=0.34, h=0.24)
    add_rect(slide, 0.82, 5.02, 5.24, 0.49, fill="FCECEC", line="E7A5A5", rounded=True)
    add_text(slide, "把未入库代码误计为入库，使入库分子虚增 60 行",
             1.00, 5.09, 4.86, 0.32, size=13.0, color=RED, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 场景二
    add_rect(slide, 6.72, 1.48, 5.65, 4.28, fill=WHITE, line=LINE, rounded=True)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.72), Inches(1.48), Inches(5.65), Inches(0.50))
    header.fill.solid(); header.fill.fore_color.rgb = rgb(BLUE); header.line.fill.background()
    add_text(slide, "场景二｜commit 数天后才入库", 6.92, 1.56, 5.22, 0.32,
             size=17, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    items = [
        ("7 月 1 日", "接受 60 行并形成本地 commit B"),
        ("7 月 10 日", "最终 commit T 入库，其中保留 AI 代码 40 行"),
        ("", "动作日与入库日不同，不能按同一天的数据直接相除"),
    ]
    for i, (left, right) in enumerate(items):
        y = 2.25 + i * 0.93
        if left:
            add_label(slide, left, 6.98, y, 1.36)
        add_text(slide, right, 8.53 if left else 7.10, y - 0.02,
                 3.42 if left else 4.80, 0.54, size=12.7,
                 valign=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
        if i < 2:
            add_chevron(slide, 9.45, y + 0.57, w=0.34, h=0.24)
    add_rect(slide, 6.98, 5.02, 5.05, 0.49, fill="FCECEC", line="E7A5A5", rounded=True)
    add_text(slide, "结果是时间批次错配：分子、分母不是同一批代码",
             7.10, 5.09, 4.80, 0.32, size=12.5, color=RED, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_rect(slide, 0.65, 5.97, 10.90, 0.50, fill=LIGHT_BLUE, line=MID_BLUE, rounded=True)
    add_text(
        slide,
        "根因：统计起点选错。必须由目标仓库的最终 commit 确认真实入库，再反推其 AI 来源；不能把本地 commit 当作最终入库事实。",
        0.86, 6.05, 10.48, 0.32, size=13.0, color=NAVY, bold=True,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
    )


def restore_protected_slide_parts(source: Path, output: Path, slide_numbers):
    with tempfile.TemporaryDirectory(prefix="git-ai-ppt-") as tmpdir:
        tmp = Path(tmpdir)
        src_dir = tmp / "src"
        out_dir = tmp / "out"
        with zipfile.ZipFile(source) as zf:
            zf.extractall(src_dir)
        with zipfile.ZipFile(output) as zf:
            zf.extractall(out_dir)

        for number in slide_numbers:
            for rel_path in [
                Path(f"ppt/slides/slide{number}.xml"),
                Path(f"ppt/slides/_rels/slide{number}.xml.rels"),
            ]:
                src = src_dir / rel_path
                dst = out_dir / rel_path
                if src.exists():
                    dst.parent.mkdir(parents=True, exist_ok=True)
                    shutil.copy2(src, dst)

        rebuilt = tmp / "restored.pptx"
        with zipfile.ZipFile(rebuilt, "w", zipfile.ZIP_DEFLATED) as zf:
            for path in sorted(out_dir.rglob("*")):
                if path.is_file():
                    zf.write(path, path.relative_to(out_dir).as_posix())
        shutil.copy2(rebuilt, output)


def build_deck():
    prs = Presentation(SOURCE)
    original = list(prs.slides)
    assert len(original) == 11
    cover, cli, attribution, win, linux, ide, profile, faq, lifecycle, problem, solution = original

    # Only the cover and the problem slide are edited among the original slides.
    title, subtitle, date = cover.shapes[0], cover.shapes[1], cover.shapes[2]
    title.left, title.top, title.width, title.height = Inches(0.84), Inches(0.52), Inches(9.2), Inches(0.80)
    set_existing_text(title, "Git AI 企业化集成与二开实践", size=32, color=WHITE, bold=True)
    subtitle.left, subtitle.top, subtitle.width, subtitle.height = Inches(1.35), Inches(2.82), Inches(10.45), Inches(0.80)
    set_existing_text(subtitle, "上游融合、Kilo v7 适配、离线交付、后台统计闭环与数据治理",
                      size=23, color=NAVY, bold=True)
    date.left, date.top, date.width, date.height = Inches(1.36), Inches(3.77), Inches(6.0), Inches(0.40)
    set_existing_text(date, "内部技术研究 · 2026 年 7 月｜基线：Git AI 1.6.16",
                      size=14.5, color=TEXT)
    redesign_problem_slide(problem)

    blank = prs.slide_layouts[6]

    # New directory.
    directory = prs.slides.add_slide(blank)
    add_title(directory, "内容目录", "原有研究页保持不动，仅按功能链路重新编排并补充二开与后台章节。")
    chapters = [
        ("01", "基础原理与边界", "CLI、插件、行级来源事实"),
        ("02", "离线部署与使用", "Windows、Linux、IDE 与 FAQ"),
        ("03", "Git AI 二开能力", "企业上报、Kilo v7、生命周期与离线交付"),
        ("04", "后台适配与统计闭环", "接入、投影、归因、明细与留存"),
        ("05", "验证与上线边界", "完成度、发布门槛与下一步"),
    ]
    for i, (num, name, desc) in enumerate(chapters):
        y = 1.54 + i * 0.93
        add_text(directory, num, 0.72, y, 0.64, 0.46, size=22, color=BLUE,
                 bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        bar = directory.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.52), Inches(y + 0.04), Inches(0.04), Inches(0.43))
        bar.fill.solid(); bar.fill.fore_color.rgb = rgb(MID_BLUE); bar.line.fill.background()
        add_text(directory, name, 1.82, y - 0.01, 3.45, 0.44, size=18,
                 color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_text(directory, desc, 5.35, y, 6.20, 0.42, size=14,
                 color=TEXT, valign=MSO_ANCHOR.MIDDLE)

    sec1 = prs.slides.add_slide(blank)
    add_section(sec1, "01", "基础原理与边界", "明确 Git AI 记录什么、插件负责什么，以及 AI / Human / Unknown 的口径。")

    sec2 = prs.slides.add_slide(blank)
    add_section(sec2, "02", "离线部署与使用", "面向内网研发环境，统一 Windows、Linux、VS Code 与 JetBrains 的安装路径。")

    sec3 = prs.slides.add_slide(blank)
    add_section(sec3, "03", "Git AI 二开能力", "企业身份、Kilo v7 采集、生命周期可靠性与可重复离线交付。")

    custom = prs.slides.add_slide(blank)
    add_title(custom, "Git AI 二开工作全景", "在上游 1.6.16 之上保留自有提交链，按四类能力维护。")
    items = [
        ("企业上报身份", "组织架构下拉与缓存\n姓名、邮箱、服务地址\nVS Code / JetBrains / CLI 共享配置"),
        ("Kilo v7 适配", "独立 preset 与安装器\n编辑、写入、补丁、命令链路\n工具、模型、运行时元数据"),
        ("生命周期可靠性", "checkpoint / commit / rewrite\nsession_event 与 Token 事实\n持久队列、重试与候选替换"),
        ("离线交付维护", "Windows x64、Linux x64 / ARM64\nVS Code / JetBrains 插件\n上游融合与可复核校验和"),
    ]
    for i, (name, body) in enumerate(items):
        x = 0.58 + (i % 2) * 6.08
        y = 1.50 + (i // 2) * 2.23
        add_card(custom, x, y, 5.72, 1.88, name, body, fill=PALE_BLUE,
                 title_size=17, body_size=13.2)
    add_text(custom, "维护原则：尽量复用上游协议与模板，通过独立适配层降低升级冲突；客户端只上报事实，最终入库归因由后台派生。",
             0.72, 5.92, 10.85, 0.38, size=13.5, color=NAVY, bold=True,
             valign=MSO_ANCHOR.MIDDLE)

    kilo = prs.slides.add_slide(blank)
    add_title(kilo, "Kilo v7：从兼容导入升级为一等采集源", "不改 Kilo 主流程，通过独立 preset、安装器与事件转换接入 Git AI。")
    add_card(kilo, 0.58, 1.48, 5.18, 4.55, "适配实现",
             "• 独立 preset：由 OpenCode 模板严格替换生成，减少上游冲突。\n\n• 多入口：覆盖 VS Code、JetBrains 与 CLI 运行时。\n\n• 文件变更：识别 edit、write、patch、bash / shell 产生的变更。\n\n• 运行元数据：上报 tool=kilo、编辑器、版本、模型、provider 与会话。\n\n• 失败边界：工具无法给出文件路径时落 Unknown；不统计 Tab 补全。",
             fill=PALE_BLUE, title_size=18, body_size=12.7)
    add_text(kilo, "统一事件链", 6.18, 1.70, 5.65, 0.36, size=18,
             color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    flow = [
        ("Kilo v7", "编辑与工具调用"),
        ("Git AI preset", "转换来源事实"),
        ("checkpoint", "记录接受行与上下文"),
        ("commit / rewrite", "生成 Notes 与生命周期事件"),
        ("后台统计", "与其他 Agent 进入同一口径"),
    ]
    for i, (left, right) in enumerate(flow):
        y = 2.28 + i * 0.68
        add_rect(kilo, 6.18, y, 5.65, 0.52, fill=WHITE, line=MID_BLUE, rounded=True)
        add_text(kilo, left, 6.36, y + 0.05, 1.72, 0.40, size=13,
                 color=DEEP_BLUE, bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_text(kilo, right, 8.17, y + 0.05, 3.42, 0.40, size=12.4,
                 color=TEXT, valign=MSO_ANCHOR.MIDDLE)

    offline = prs.slides.add_slide(blank)
    add_title(offline, "离线交付与上游维护：把二开做成可重复发布",
              "当前自有分支已融合 upstream/main 1.6.16，并产出可复核离线包。")
    stages = [
        ("上游基线", "1.6.16\n融合 53 个上游提交"),
        ("冲突融合", "rewrite 流式解析\n保留 hunks 采集"),
        ("跨平台构建", "Windows x64\nLinux x64 / ARM64"),
        ("插件产物", "VS Code VSIX\nJetBrains ZIP"),
        ("离线校验", "根目录 SHA256SUMS\n包内校验和"),
    ]
    for i, (name, body) in enumerate(stages):
        x = 0.55 + i * 2.48
        add_step(offline, i + 1, name, body, x, 1.62, 2.18, 3.45)
        if i < 4:
            add_chevron(offline, x + 2.21, 3.18)
    add_rect(offline, 0.72, 5.42, 10.90, 0.70, fill=LIGHT_BLUE, line=MID_BLUE, rounded=True)
    add_text(offline, "当前结果：离线包校验和已通过。正式发包前仍需完成真实 Windows x64 安装、Hook 与 commit 冒烟。",
             0.94, 5.56, 10.46, 0.38, size=13.4, color=NAVY, bold=True,
             valign=MSO_ANCHOR.MIDDLE)

    sec4 = prs.slides.add_slide(blank)
    add_section(sec4, "04", "后台适配与统计闭环", "后台负责幂等接入、事件投影、最终提交确认、统一指标、事件明细和数据留存。")

    backend = prs.slides.add_slide(blank)
    add_title(backend, "后台适配全链路", "客户端上报可追溯事实；后台把事实转换为可查询、可确认、可重建的统计结果。")
    stages = [
        ("Git AI / Kilo", "checkpoint\ncommit / rewrite\nsession / Token"),
        ("批量接入", "批次与事件幂等\n失败可重放\n保留处理状态"),
        ("事件规范化", "header / payload 分离\n组织、用户、仓库\n工具、模型、会话"),
        ("归因与生命周期", "候选建立 / 替换\n目标分支确认\n行级去重"),
        ("统计与管理端", "日汇总与统一口径\n事件明细\n提交定位与排障"),
    ]
    for i, (name, body) in enumerate(stages):
        x = 0.48 + i * 2.54
        add_card(backend, x, 1.58, 2.18, 3.86, name, body, fill=PALE_BLUE,
                 title_size=13.7, body_size=12.1)
        if i < 4:
            add_chevron(backend, x + 2.24, 3.18, w=0.25, h=0.38)
    add_text(backend, "隐私边界：后台只展示统计与事件事实，不展示原始 prompt、回复或 transcript；客户端不承担 committed 权威归因。",
             0.72, 5.72, 10.90, 0.42, size=13.6, color=NAVY, bold=True,
             valign=MSO_ANCHOR.MIDDLE)

    inventory = prs.slides.add_slide(blank)
    add_title(inventory, "为了适配 Git AI，后台新增了哪些能力")
    rows = [
        ("协议与入库", "批量接收 Git AI / Kilo 事件；按事件键幂等；header 与 payload 分离；支持失败重放。"),
        ("事件投影", "checkpoint 转接受事实；commit / rewrite 转候选与确认；session、Token 统一规范化。"),
        ("统计归因", "目标分支最终 commit 确认真实入库；行级去重；生命周期替换；日汇总可由规范账本重建。"),
        ("现行产品口径", "Git AI 生成行数按接受行数展示，接受率固定为 100%；入库确认仍以目标仓库事实为准。"),
        ("管理端页面", "左侧拆分“Git AI 明细”和“Kilo 提交明细”；支持用户、模块、事件、状态与 commit 定位。"),
        ("运维治理", "30 天热明细查询；会话上下文与 Token 水位长期固化；清理任务先预览，物理删除默认关闭。"),
    ]
    add_table(inventory, 0.58, 1.35, 11.80, 4.92, ["能力域", "后台适配内容"], rows, [2.15, 9.65])

    retention = prs.slides.add_slide(blank)
    add_title(retention, "事件明细与 30 天留存策略",
              "事件页用于定位每次提交和上报问题；留存按“热明细可见、规范账本长期可重建”分层。")
    add_card(retention, 0.58, 1.48, 5.72, 4.60, "事件明细展示",
             "• checkpoint：接受行、工具、模型、会话、文件与投影状态。\n\n• commit / rewrite：提交、hunk、候选替换、确认与业务状态。\n\n• session_event：运行时、编辑器、版本、provider 与生命周期。\n\n• Token：只保留计量、水位与去重事实，不展示原始对话。\n\n• 两个视图：Git AI 明细与 Kilo 提交明细分别查询。",
             fill=WHITE, title_size=18, body_size=12.8)
    add_rect(retention, 6.65, 1.48, 5.72, 4.60, fill=PALE_BLUE, line=MID_BLUE, rounded=True)
    add_text(retention, "数据分层", 6.95, 1.74, 5.02, 0.36, size=18,
             color=NAVY, bold=True)
    layers = [
        ("0–30 天", "热明细", "页面和接口可查询完整事件事实。", GREEN),
        (">30 天", "长期规范账本", "保留幂等、重放、会话上下文、Token 水位与重建依据。", DEEP_BLUE),
        ("物理清理", "尚未开启", "先做 dry-run、归档回放和汇总重建一致性验收。", AMBER),
    ]
    for i, (period, name, body, color) in enumerate(layers):
        y = 2.35 + i * 1.02
        add_label(retention, period, 6.95, y, 1.14, fill=color, color=WHITE)
        add_text(retention, name, 8.28, y - 0.01, 1.50, 0.32, size=14.2,
                 color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_text(retention, body, 9.72, y - 0.03, 2.18, 0.60, size=11.8,
                 color=TEXT, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
    add_text(retention, "数据库增长的主要原因之一，是历史事件 payload 和保障统计正确性的账本尚未物理清理。",
             6.95, 5.43, 4.88, 0.42, size=12.8, color=RED, bold=True)

    sec5 = prs.slides.add_slide(blank)
    add_section(sec5, "05", "验证与上线边界", "完成度以可验证证据为准：合并、构建、回归、真实运行与数据清理分别设门槛。")

    release = prs.slides.add_slide(blank)
    add_title(release, "当前完成度与正式发包前清单")
    add_card(release, 0.58, 1.35, 5.70, 4.92, "已经完成并可复核",
             "• Git AI 上游已融合到 1.6.16，自有二开保留在独立提交链。\n\n• rewrite 冲突已手工融合，Kilo hunks 与流式解析同时保留。\n\n• Windows / Linux / VS Code / JetBrains 离线产物已生成，校验和通过。\n\n• Git AI 与 Kilo 事件已进入统一接入、投影、归因和管理端明细链路。\n\n• 后台现行生成行数、接受率与目标 commit 入库口径已统一。",
             header=GREEN, fill=PALE_BLUE, title_size=17, body_size=13.0)
    add_card(release, 6.68, 1.35, 5.70, 4.92, "正式发包前必须收口",
             "1. 在真实 Windows x64 完成 CLI 安装、Hook、checkpoint 与 commit 冒烟。\n\n2. 在 VS Code、JetBrains 与 Kilo v7 各跑一条真实编辑到后台明细的闭环。\n\n3. 复核匿名接入、组织身份绑定与网关 / Token 边界。\n\n4. 对 retention 执行 dry-run、归档回放和汇总重建一致性验收。\n\n5. 重新生成正式离线包，记录版本、校验和与回归结论。",
             header=AMBER, fill=WHITE, title_size=17, body_size=13.0)

    # Reorder only. Original slides 2-9 and 11 keep their slide parts and internal structure.
    desired = [
        cover, directory, sec1, cli, attribution, sec2, win, linux, ide, faq,
        sec3, custom, profile, kilo, lifecycle, offline,
        sec4, backend, inventory, retention, problem, solution,
        sec5, release,
    ]
    current_slides = list(prs.slides)
    id_map = {id(slide): sld_id for slide, sld_id in zip(current_slides, list(prs.slides._sldIdLst))}
    for sld_id in list(prs.slides._sldIdLst):
        prs.slides._sldIdLst.remove(sld_id)
    for slide in desired:
        prs.slides._sldIdLst.append(id_map[id(slide)])

    prs.save(OUTPUT)

    # python-pptx serializes every slide; restore protected original slide XML byte-for-byte.
    restore_protected_slide_parts(SOURCE, OUTPUT, slide_numbers=[2, 3, 4, 5, 6, 7, 8, 9, 11])
    print(f"saved: {OUTPUT}")
    print(f"slides: {len(desired)}")


if __name__ == "__main__":
    build_deck()
